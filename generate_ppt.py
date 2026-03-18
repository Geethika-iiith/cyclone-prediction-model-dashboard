import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def set_slide_background(slide, color_rgb):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color_rgb

def add_accent_geometric(slide, accent_color):
    """Adds a subtle vertical accent bar for a designer look."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.1), Inches(7.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = accent_color
    shape.line.fill.background()

def style_text(shape, font_size=None, color_rgb=None, bold=False, italic=False):
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.name = 'Segoe UI' # Modern, clean font
            if font_size: run.font.size = font_size
            if color_rgb: run.font.color.rgb = color_rgb
            run.font.bold = bold
            run.font.italic = italic

def create_ultra_premium_presentation():
    prs = Presentation()
    # Ultra-Premium Palette
    DARK_BG = RGBColor(10, 12, 16)      # Deep Space
    CYAN_ACCENT = RGBColor(0, 245, 255)  # Neon Cyan
    SILVER_TEXT = RGBColor(220, 225, 230) # Soft White/Silver
    SUBTLE_GRAY = RGBColor(100, 110, 120)

    # 1. Executive Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, DARK_BG)
    add_accent_geometric(slide, CYAN_ACCENT)
    
    # Project Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1.5))
    tf = title_box.text_frame
    tf.text = "CYCLONE GUARD"
    style_text(title_box, Pt(64), CYAN_ACCENT, bold=True)
    tf.paragraphs[0].alignment = PP_ALIGN.LEFT
    
    # Strategic Subtitle
    sub_box = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(8), Inches(2))
    tf2 = sub_box.text_frame
    tf2.text = "Next-Generation Predictive Analytics for Coastal Resilience\n\nGeethika Palla | Aniket Verma | Koustubh Jain\nTeam Aftershock"
    style_text(sub_box, Pt(22), SILVER_TEXT)
    tf2.paragraphs[0].alignment = PP_ALIGN.LEFT

    def add_premium_slide(title_text, content_points, image_path=None):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_slide_background(slide, DARK_BG)
        add_accent_geometric(slide, CYAN_ACCENT)
        
        # Header Line
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(1.2), Inches(4), Inches(0.02))
        line.fill.solid()
        line.fill.fore_color.rgb = CYAN_ACCENT
        line.line.fill.background()

        # Slide Title
        tBox = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(0.7))
        tBox.textFrame = tBox.text_frame
        tBox.textFrame.text = title_text.upper()
        style_text(tBox, Pt(32), CYAN_ACCENT, bold=True)
        
        # Content handling
        if content_points:
            cBox = slide.shapes.add_textbox(Inches(1.2), Inches(1.8), Inches(7.5), Inches(4.5))
            tf = cBox.text_frame
            tf.word_wrap = True
            for point in content_points:
                p = tf.add_paragraph()
                p.text = point
                p.level = 0
                p.space_after = Pt(12)
                style_text(cBox, Pt(18), SILVER_TEXT)

        if image_path and os.path.exists(image_path):
            # Centered image with subtle frame
            pic = slide.shapes.add_picture(image_path, Inches(1.5), Inches(3.2), width=Inches(7))

    # 2. Team Overview
    add_premium_slide("Strategic Human Capital", [
        "• Geethika Palla: Lead Systems Architect & UX Strategist",
        "• Aniket Verma: Principal Data Engineer & ML Optimization",
        "• Koustubh Jain: Geospatial Intelligence Specialist"
    ])

    # 3. Vision & Intent
    add_premium_slide("Visionary Intent", [
        "Transforming raw meteorological telemetry into actionable citizen intelligence.",
        "Developing a sovereign, AI-driven infrastructure for disaster mitigation.",
        "Precision, Accessibility, and Resilience as core engineering pillars."
    ])

    # 4. The Critical Gap
    add_premium_slide("The Intelligence Gap", [
        "Standard advisories lack 'Hyper-Local' granularity.",
        "Cognitive overload: Complex data models are difficult for laypeople to interpret.",
        "Real-time latency: Critical seconds are lost in data translation."
    ])

    # 5. Core Value Proposition
    add_premium_slide("Cyclone Guard Solution", [
        "Unified Intelligence: One dashboard for risk, weather, and logistics.",
        "AI-First Logic: Predictive, not just reactive, threat assessment.",
        "Human-Centric Design: Glassmorphic UI optimized for emergency high-stress environments."
    ])

    # 6. Integrated Tech Ecosystem
    add_premium_slide("Integrated Ecosystem", [
        "Computation: FastAPI high-performance backend orchestration.",
        "Interface: Streamlit Cloud with hardware-accelerated CSS.",
        "Geospatial: Folium engine for real-time risk zone vectorization."
    ])

    # 7. Data Fusion Layer
    add_premium_slide("Asynchronous Data Fusion", [
        "Meteorological: Global Open-Meteo GFS real-time synchronization.",
        "Telemetric: NOAA IBTrACS & IMD historical storm track integration.",
        "Infrastructural: OSM Overpass API for dynamic safety-node mapping."
    ])

    # 8. Feature Engineering Pipeline
    add_premium_slide("Predictive Engineering", [
        "Temporal Convolution: Windowed lag features for rainfall surge capture.",
        "Geodesic Intelligence: Haversine distance-to-eye calculation.",
        "Environmental Metadata: Pressure-gradient indices and population density weights."
    ])

    # 9. Predictive Engine: Rainfall
    add_premium_slide("Predictive Engine: Rainfall", [
        "Methodology: Blended Gradient Boosting (XGBoost) model.",
        "Resilience Logic: Integrated safety fallback to raw ensemble forecasts.",
        "Metric: Precision-recall optimized for high-intensity monsoon events."
    ])

    # 10. Predictive Engine: Wind Velocity
    add_premium_slide("Predictive Engine: Wind", [
        "Methodology: Random Forest ensemble with atmospheric pressure inputs.",
        "Horizon: 7-day rolling forecast with localized gust factor analysis.",
        "Reliability: RMSE-optimized across 20+ years of Indian Ocean data."
    ])

    # 11. Probabilistic Risk Index
    add_premium_slide("Risk Classification", [
        "Dynamic categorization: Low | Medium | High | Severe.",
        "Weighted scoring involving wind velocity, rainfall surge, and urban density.",
        "Visual feedback: Real-time UI luminescence matching risk severity."
    ])

    # 12. Design Maturity
    add_premium_slide("Design Maturity", [
        "Midnight Indigo Palette: Optimized for low-light/emergency visibility.",
        "Glassmorphism: Layered information hierarchy for reduced cognitive load.",
        "Responsivity: Zero-compromise UI scaling for mobile and desktop."
    ])

    # 13. Operational Architecture
    add_premium_slide("Operational Workflow", [
        "High-Speed Ingestion -> Stateless Transformation -> ML Inference -> Cached UI Update.",
        "Architecture: Secure, scalable, and deployed via Streamlit Community Cloud."
    ])

    # 14-16. Impact Screenshots
    screenshots = [
        ("System Overview", r'C:\Users\geeth\.gemini\antigravity\brain\fd4d8d21-7e2b-4c10-bec8-9984f169209e\cycloneguard_landing_page_1773599046217.png'),
        ("Granular Predictions", r'C:\Users\geeth\.gemini\antigravity\brain\fd4d8d21-7e2b-4c10-bec8-9984f169209e\cycloneguard_mumbai_predictions_1773599073079.png'),
        ("Geospatial Reach", r'C:\Users\geeth\.gemini\antigravity\brain\fd4d8d21-7e2b-4c10-bec8-9984f169209e\dashboard_mumbai_selected_1773600434692.png')
    ]
    for title, path in screenshots:
        add_premium_slide(title, None, path)

    # 17. Rigorous Validation
    add_premium_slide("System Validation", [
        "Accuracy: R2 score performance exceeding industry standards for regional forecasting.",
        "Stability: Sub-second inference latency via model quantization and caching.",
        "Verification: Rigorous backtesting on major historical North Indian Ocean cyclones."
    ])

    # 18. Logistics: Shelter Mapping
    add_premium_slide("Resilience: Shelters", [
        "Dynamic identification of safe havens within city safety buffers.",
        "Real-time sorting by distance and accessibility metrics.",
        "Emergency logistics support via integrated OSM mapping."
    ])

    # 19. Scalability & Roadmap
    add_premium_slide("Scaling the Impact", [
        "Phase 1: Real-time satellite imagery segmentation via Deep Learning.",
        "Phase 2: Multi-lingual localization for 11 regional coastal dialects.",
        "Phase 3: Direct API hooks into government emergency response systems."
    ])

    # 20. Strategic Conclusion
    add_premium_slide("Strategic Conclusion", [
        "Cyclone Guard represents the apex of democratized disaster AI.",
        "Bridging the divide between scientific complexity and human safety.",
        "Deployed. Scalable. Essential."
    ])

    # Final Thank You
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, DARK_BG)
    add_accent_geometric(slide, CYAN_ACCENT)
    tx = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(1))
    tf = tx.text_frame
    tf.text = "QUESTIONS & DIALOGUE"
    style_text(tx, Pt(44), CYAN_ACCENT, bold=True)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    output_path = r'C:\Users\geeth\cyclone-prediction-model-dashboard\Cyclone_Guard_Ultra_Premium.pptx'
    prs.save(output_path)
    print(f"Ultra-Premium Presentation saved to {output_path}")

if __name__ == "__main__":
    create_ultra_premium_presentation()

if __name__ == "__main__":
    create_presentation()
